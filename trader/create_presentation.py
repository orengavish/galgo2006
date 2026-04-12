#!/usr/bin/env python3
"""
Galao System — Presentation Generator
Creates galao_presentation.pptx for developer partner briefing.

Usage:
    python create_presentation.py
    python create_presentation.py --self-test
"""

import sys
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG    = RGBColor(0x1A, 0x23, 0x3A)   # dark navy
PANEL = RGBColor(0x24, 0x34, 0x50)   # slightly lighter panel
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xFF, 0xA5, 0x00)   # accent
BLUE  = RGBColor(0x5B, 0x9B, 0xD5)   # secondary
GRAY  = RGBColor(0xBB, 0xBB, 0xBB)   # subtext
GREEN = RGBColor(0x70, 0xAD, 0x47)   # profit / pass
RED   = RGBColor(0xFF, 0x5C, 0x5C)   # loss / fail
YELLOW= RGBColor(0xFF, 0xD7, 0x00)   # warning

SW = Inches(13.33)   # slide width  (16:9)
SH = Inches(7.5)     # slide height

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def tb(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
       align=PP_ALIGN.LEFT, italic=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.color.rgb = color
    r.font.bold  = bold
    r.font.italic = italic
    return shape

def mtb(slide, lines, x, y, w, h, default_size=16, default_color=WHITE,
        default_bold=False, align=PP_ALIGN.LEFT):
    """Multi-line text box. Each line: str or (text, color, bold, size)."""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = shape.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if isinstance(line, tuple):
            txt  = line[0]
            col  = line[1] if len(line) > 1 else default_color
            bld  = line[2] if len(line) > 2 else default_bold
            sz   = line[3] if len(line) > 3 else default_size
        else:
            txt, col, bld, sz = line, default_color, default_bold, default_size
        r = p.add_run()
        r.text = str(txt)
        r.font.size  = Pt(sz)
        r.font.color.rgb = col
        r.font.bold  = bld
    return shape

def box(slide, x, y, w, h, fill=PANEL, line=None, line_w=1):
    from pptx.util import Pt as Pt2
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt2(line_w)
    else:
        s.line.fill.background()
    return s

def hline(slide, y, color=GOLD, x=0.5, w=12.33, thickness=2):
    from pptx.util import Pt as Pt2
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt2(thickness))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def header(slide, title, subtitle=None):
    tb(slide, title, 0.5, 0.25, 12.33, 0.75, size=34, bold=True)
    hline(slide, 1.1, GOLD)
    if subtitle:
        tb(slide, subtitle, 0.5, 1.15, 12.33, 0.45, size=15, color=GRAY)

def lifecycle_bar(slide, current=0):
    """Progress bar showing lifecycle step. current=0 means no highlight."""
    steps = ["Day Start","DB Write","IB Submit","Fill","OCO Active","Exit","Replenish","Shutdown"]
    x = 0.3
    w = 12.73 / len(steps)
    for i, step in enumerate(steps):
        fill = GOLD if i+1 == current else RGBColor(0x30, 0x45, 0x60)
        box(slide, x + i*w, 7.0, w-0.05, 0.35, fill=fill)
        col = BG if i+1 == current else GRAY
        tb(slide, step, x+i*w, 7.0, w-0.05, 0.35,
           size=9, color=col, bold=(i+1==current), align=PP_ALIGN.CENTER)

def tbl(slide, headers, rows, x, y, w, h):
    t = slide.shapes.add_table(len(rows)+1, len(headers),
                                Inches(x), Inches(y), Inches(w), Inches(h)).table
    # header row
    for ci, hdr in enumerate(headers):
        c = t.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = GOLD
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr; r.font.bold = True
        r.font.size = Pt(12); r.font.color.rgb = BG
    # data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = t.cell(ri+1, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = PANEL if ri%2==0 else BG
            p = c.text_frame.paragraphs[0]
            r = p.add_run()
            if isinstance(val, tuple):
                r.text = val[0]; r.font.color.rgb = val[1]
            else:
                r.text = str(val); r.font.color.rgb = WHITE
            r.font.size = Pt(11)
    return t

# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title_page(prs):
    s = blank(prs)
    box(s, 0, 0, 13.33, 7.5, fill=BG)
    # accent bar left
    box(s, 0, 0, 0.15, 7.5, fill=GOLD)
    tb(s, "GALAO", 0.4, 1.2, 12.0, 1.2, size=72, bold=True, color=WHITE)
    tb(s, "Automated Intraday Futures Scalper", 0.4, 2.5, 12.0, 0.6,
       size=26, color=GOLD, bold=True)
    hline(s, 3.3, GOLD, x=0.4, w=10.0)
    tb(s, "Developer Partner Briefing  ·  MES Futures  ·  April 2026",
       0.4, 3.5, 12.0, 0.5, size=16, color=GRAY)
    tb(s, "Interactive Brokers  ·  ib_insync  ·  Python  ·  SQLite",
       0.4, 4.1, 12.0, 0.4, size=14, color=BLUE)

def slide_agenda(prs):
    s = blank(prs)
    header(s, "Agenda")
    items = [
        ("1.  The Thesis — Why Critical Lines Work",       GOLD,  True,  20),
        ("",                                                WHITE, False, 8),
        ("2.  Strategy Overview — What We Trade",          WHITE, False, 18),
        ("",                                                WHITE, False, 8),
        ("3.  Order Logic — Toggle Rule & OCO Brackets",   WHITE, False, 18),
        ("",                                                WHITE, False, 8),
        ("4.  Trade Lifecycle  (main section)",             GOLD,  True,  20),
        ("     Day Start  →  DB Write  →  IB Submit  →  Fill  →  Exit  →  Replenish  →  Shutdown",
                                                            BLUE,  False, 14),
        ("",                                                WHITE, False, 8),
        ("5.  Architecture & Database Schema",              WHITE, False, 18),
        ("",                                                WHITE, False, 8),
        ("6.  A/B Testing Framework",                       WHITE, False, 18),
        ("",                                                WHITE, False, 8),
        ("7.  Dev Quality — Self-test, Regression, Versioning", WHITE, False, 18),
    ]
    mtb(s, items, 1.0, 1.3, 11.0, 5.8, align=PP_ALIGN.LEFT)

def slide_thesis(prs):
    s = blank(prs)
    header(s, "The Thesis — Critical Lines", "Why certain price levels behave differently")

    box(s, 0.4, 1.4, 5.8, 5.6, fill=PANEL, line=GOLD)
    tb(s, "What is a Critical Line?", 0.6, 1.5, 5.4, 0.5, size=18, bold=True, color=GOLD)
    mtb(s, [
        "A price level where the market shows",
        "significantly higher activity than surrounding prices.",
        "",
        ("→  Hit 20× per day  vs  avg 5× for normal prices", BLUE, True, 16),
        "",
        "Caused by high-volume clustering at that price —",
        "institutions, stops, and liquidity sitting there.",
        "",
        "The market is drawn back to these levels",
        "repeatedly throughout the session.",
    ], 0.6, 2.1, 5.4, 4.6, default_size=15)

    box(s, 6.5, 1.4, 6.4, 5.6, fill=PANEL, line=BLUE)
    tb(s, "The Vibration Effect", 6.7, 1.5, 6.0, 0.5, size=18, bold=True, color=GOLD)
    mtb(s, [
        "Price oscillates around critical lines.",
        "Every touch is a trading opportunity.",
        "",
        ("Price above 6250:", BLUE, True, 15),
        "  →  drops to 6250  →  bounces up",
        ("Price below 6250:", BLUE, True, 15),
        "  →  rises to 6250  →  bounces down",
        "",
        "Strategy: place orders at the line,",
        "harvest the bounce, repeat on replenishment.",
        "",
        ("Strength 1-3 labels the line's reliability",
         YELLOW, False, 13),
        ("for A/B testing.",
         YELLOW, False, 13),
    ], 6.7, 2.1, 6.0, 4.6, default_size=15)

def slide_strategy(prs):
    s = blank(prs)
    header(s, "Strategy Overview", "What we trade and how")

    items = [
        [("Instrument",    GOLD, True,  14), ("MES — Micro E-mini S&P 500 Futures (CME)",  WHITE, False, 14)],
        [("Why Micro",     GOLD, True,  14), ("$5/pt cost of learning. Allows 50+ simultaneous positions.", WHITE, False, 14)],
        [("Direction",     GOLD, True,  14), ("Both LONG and SHORT at every critical line.",WHITE, False, 14)],
        [("Entry type",    GOLD, True,  14), ("Limit orders (LMT) — bounce strategy, not breakout.",WHITE, False, 14)],
        [("Brackets",      GOLD, True,  14), ("Symmetric TP = SL. Sizes from config: [2, 4, 8, 16] pts.", WHITE, False, 14)],
        [("Trade length",  GOLD, True,  14), ("Short by design. Stagnation exit at 60s if no movement.", WHITE, False, 14)],
        [("Session",       GOLD, True,  14), ("CME open +30min  →  CME close -60min",       WHITE, False, 14)],
        [("Port",          GOLD, True,  14), ("PAPER port only. LIVE port = data only.",     WHITE, False, 14)],
        [("Phase",         GOLD, True,  14), ("Learning & A/B testing. P&L optimisation later.", YELLOW, False, 14)],
    ]
    y = 1.4
    for pair in items:
        box(s, 0.4, y, 2.6, 0.36, fill=PANEL)
        box(s, 3.1, y, 9.8, 0.36, fill=BG)
        tb(s, pair[0][0], 0.5, y+0.04, 2.4, 0.3, size=pair[0][3], color=pair[0][1], bold=pair[0][2])
        tb(s, pair[1][0], 3.2, y+0.04, 9.6, 0.3, size=pair[1][3], color=pair[1][1], bold=pair[1][2])
        y += 0.44

def slide_critical_lines(prs):
    s = blank(prs)
    header(s, "Critical Lines — Input Format", "Manual input for V1. Auto-detection planned for V2.")

    box(s, 0.4, 1.4, 5.8, 2.2, fill=PANEL, line=GOLD)
    tb(s, "File Format", 0.6, 1.5, 5.4, 0.4, size=16, bold=True, color=GOLD)
    tb(s, "data/critical_lines/levels_daily_YYYYMMDD.txt",
       0.6, 2.0, 5.4, 0.35, size=13, color=BLUE, italic=True)
    mtb(s, [
        ("SYMBOL,  PRICE,    STRENGTH", YELLOW, True, 14),
        ("MES,     6250.00,  3",         GREEN,  False,14),
        ("MES,     6300.00,  2",         GREEN,  False,14),
        ("MES,     6180.00,  1",         GREEN,  False,14),
        ("MES,     6400.00,  3",         GREEN,  False,14),
    ], 0.6, 2.45, 5.4, 1.0, default_size=14)

    box(s, 0.4, 3.8, 5.8, 3.1, fill=PANEL, line=BLUE)
    tb(s, "Strength Scale", 0.6, 3.9, 5.4, 0.4, size=16, bold=True, color=GOLD)
    tbl(s,
        ["Strength", "Meaning"],
        [
            [("1", YELLOW), "Weak — low confidence, fewer historical hits"],
            [("2", BLUE),   "Medium — reliable, consistent touches"],
            [("3", GREEN),  "Strong — very high confidence, high volume node"],
        ],
        0.5, 4.4, 5.7, 2.2
    )

    box(s, 6.5, 1.4, 6.4, 5.5, fill=PANEL, line=BLUE)
    tb(s, "Rules", 6.7, 1.5, 6.0, 0.4, size=16, bold=True, color=GOLD)
    mtb(s, [
        "•  Up to 10 lines per symbol per day",
        "",
        "•  Strength is recorded on every order",
        "   for A/B analysis — does NOT affect",
        "   order logic in V1",
        "",
        "•  Same file can be reused for a week",
        "",
        "•  Missing file → that symbol is blocked",
        "   from trading that day",
        "",
        ("•  V2 goal: auto-detect lines from", YELLOW, False, 14),
        ("   volume profile in Fetcher CSV data", YELLOW, False, 14),
    ], 6.7, 2.0, 6.0, 4.7, default_size=15)

def slide_toggle(prs):
    s = blank(prs)
    header(s, "Order Logic — The Toggle Rule",
           "Entry order type depends on which side of the line price is currently on")

    box(s, 0.4, 1.4, 5.8, 2.6, fill=PANEL, line=GOLD)
    tb(s, "Why toggle?", 0.6, 1.5, 5.4, 0.4, size=16, bold=True, color=GOLD)
    mtb(s, [
        "LMT BUY at 6250 when price is at 6200",
        "(below the line) would fill IMMEDIATELY —",
        "price is already below 6250.",
        "",
        ("The toggle prevents this.", RED, True, 15),
        "Order type is chosen based on which side",
        "price is currently on.",
        "",
        "Re-evaluated on every replenishment.",
    ], 0.6, 2.0, 5.4, 3.6, default_size=14)

    tbl(s,
        ["Current Price", "BUY Entry", "SELL Entry", "Logic"],
        [
            [("Price ABOVE line", GREEN), ("LMT BUY",  BLUE),   ("STP SELL", YELLOW),
             "Catch price dropping to line"],
            [("Price BELOW line", RED),   ("STP BUY",  YELLOW), ("LMT SELL", BLUE),
             "Catch price rising to line"],
        ],
        6.4, 1.4, 6.5, 1.8
    )

    box(s, 6.4, 3.4, 6.5, 3.5, fill=PANEL, line=BLUE)
    tb(s, "Example  —  Line at 6250,  Current price 6280  (above)",
       6.6, 3.5, 6.1, 0.4, size=14, bold=True, color=GOLD)
    mtb(s, [
        ("LMT BUY  @ 6250", BLUE,   True,  16),
        "  TP = 6252  |  SL = 6248  (bracket 2pt)",
        "",
        ("STP SELL @ 6250", YELLOW, True,  16),
        "  TP = 6248  |  SL = 6252  (bracket 2pt)",
        "",
        "Both trigger when price drops to 6250.",
        ("Net IB position = 0  (intentional — virtual legs)", GRAY, False, 13),
    ], 6.6, 4.0, 6.1, 2.7, default_size=14)

def slide_bracket(prs):
    s = blank(prs)
    header(s, "OCO Bracket Structure", "Every entry order has exactly two children: TP and SL")

    # Bracket diagram
    box(s, 0.4, 1.4, 6.0, 5.5, fill=PANEL, line=GOLD)
    tb(s, "Bracket Anatomy", 0.6, 1.5, 5.6, 0.4, size=16, bold=True, color=GOLD)

    # draw bracket
    box(s, 1.5, 2.1, 3.8, 0.55, fill=RGBColor(0x30,0x50,0x80), line=BLUE)
    tb(s, "PARENT  —  LMT BUY @ 6250", 1.6, 2.18, 3.6, 0.38,
       size=14, bold=True, color=WHITE)

    # TP box
    box(s, 0.7, 3.2, 2.2, 0.55, fill=RGBColor(0x1A,0x50,0x2A), line=GREEN)
    tb(s, "TP  LMT SELL @ 6252", 0.8, 3.28, 2.0, 0.38, size=13, color=GREEN, bold=True)

    # SL box
    box(s, 3.3, 3.2, 2.2, 0.55, fill=RGBColor(0x50,0x1A,0x1A), line=RED)
    tb(s, "SL  STP SELL @ 6248", 3.4, 3.28, 2.0, 0.38, size=13, color=RED, bold=True)

    tb(s, "OCO: when one fills → other cancels automatically",
       0.6, 4.0, 5.4, 0.4, size=13, color=YELLOW, italic=True)

    tb(s, "transmit=False   transmit=False   transmit=True",
       0.6, 4.5, 5.4, 0.35, size=12, color=GRAY, italic=True)
    tb(s, "← last order transmits entire bracket to IB",
       0.6, 4.85, 5.4, 0.35, size=12, color=GRAY, italic=True)

    # Bracket sizes panel
    box(s, 6.5, 1.4, 6.4, 5.5, fill=PANEL, line=BLUE)
    tb(s, "A/B Test — Bracket Sizes", 6.7, 1.5, 6.0, 0.4, size=16, bold=True, color=GOLD)
    tbl(s,
        ["Size (pts)", "TP", "SL", "Profile"],
        [
            ["2",  "+2", "-2", ("High freq, small profit", GREEN)],
            ["4",  "+4", "-4", "Medium frequency"],
            ["8",  "+8", "-8", "Lower frequency"],
            ["16", "+16","-16",("Rare, large profit",     YELLOW)],
        ],
        6.6, 2.05, 6.0, 2.4
    )
    mtb(s, [
        ("Active brackets set in config.yaml", BLUE, False, 13),
        "  orders.active_brackets: [2, 4]",
        "",
        "All bracket sizes run simultaneously.",
        "Analyzer will compare win rates per size.",
        "",
        ("Brackets are always symmetric.", YELLOW, True, 14),
        "TP distance = SL distance from entry.",
    ], 6.7, 4.6, 6.0, 2.2, default_size=14)

def slide_lifecycle_overview(prs):
    s = blank(prs)
    header(s, "Trade Lifecycle — Overview", "From critical line to closed position")

    steps = [
        ("1\nDAY\nSTART",   "Decider reads\ncritical lines\nfile",         GOLD),
        ("2\nDB\nWRITE",    "Commands\nwritten\nPENDING",                  BLUE),
        ("3\nIB\nSUBMIT",   "Broker builds\nOCO bracket\n→ PAPER 4002",   BLUE),
        ("4\nFILL",         "Price hits\ncritical line\nAll 4 fill",       GREEN),
        ("5\nOCO\nACTIVE",  "TP + SL\nchildren\nactivated",               GREEN),
        ("6\nEXIT",         "TP / SL /\nStagnation\nexits",               YELLOW),
        ("7\nREPLENISH",    "Decider writes\nnew PENDING\ncmd",           GOLD),
        ("8\nSHUTDOWN",     "T-60min:\ncancel all,\nexit all",            RED),
    ]

    bw = 1.4
    bh = 3.8
    bx = 0.35
    by = 2.0

    for i, (title, desc, color) in enumerate(steps):
        x = bx + i * (bw + 0.12)
        box(s, x, by, bw, bh, fill=PANEL, line=color)
        tb(s, title, x+0.05, by+0.1, bw-0.1, 0.85,
           size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        hline(s, by+1.0, color, x=x+0.1, w=bw-0.2, thickness=1)
        tb(s, desc, x+0.05, by+1.1, bw-0.1, bh-1.2,
           size=12, color=GRAY, align=PP_ALIGN.CENTER)

        if i < len(steps)-1:
            tb(s, "→", x+bw+0.01, by+1.6, 0.15, 0.4,
               size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

    tb(s, "Replenishment cycles continuously until Shutdown — each fill = new set of orders at same line",
       0.5, 6.1, 12.33, 0.4, size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

def slide_lc_daystart(prs):
    s = blank(prs)
    header(s, "Lifecycle 1 — Day Start", "Decider reads lines and generates all commands")
    lifecycle_bar(s, 1)

    box(s, 0.4, 1.4, 5.7, 5.4, fill=PANEL, line=GOLD)
    tb(s, "Decider Actions  (09:00 CT)", 0.6, 1.5, 5.3, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        "1.  Read levels_daily_20260407.txt",
        "2.  Fetch current price from DB",
        "    →  MES current price = 6280",
        "3.  Evaluate toggle for each line:",
        ("    6280 > 6250  →  LMT BUY + STP SELL", BLUE, True, 14),
        "4.  Round prices to 0.25 tick",
        "5.  Write 4 commands to DB (2 brackets × 2 dirs)",
        "6.  Set session_state = RUNNING",
    ], 0.6, 2.0, 5.3, 4.6, default_size=14)

    box(s, 6.3, 1.4, 6.6, 5.4, fill=PANEL, line=BLUE)
    tb(s, "DB  commands  table  (after write)", 6.5, 1.5, 6.2, 0.4,
       size=15, bold=True, color=GOLD)
    tbl(s,
        ["id","dir","type","entry","TP","SL","bkt","str","status"],
        [
            ["1",("BUY", GREEN),("LMT",BLUE),"6250","6252","6248","2","2",("PENDING",YELLOW)],
            ["2",("BUY", GREEN),("LMT",BLUE),"6250","6254","6246","4","2",("PENDING",YELLOW)],
            ["3",("SELL",RED),  ("STP",YELLOW),"6250","6248","6252","2","2",("PENDING",YELLOW)],
            ["4",("SELL",RED),  ("STP",YELLOW),"6250","6246","6254","4","2",("PENDING",YELLOW)],
        ],
        6.4, 2.05, 6.4, 2.4
    )
    mtb(s, [
        ("4 commands = 2 directions × 2 bracket sizes", GOLD, True, 14),
        "",
        "Each command will become one OCO bracket in IB.",
        "Decider's job ends here until fills come back.",
    ], 6.5, 4.6, 6.2, 2.0, default_size=14)

def slide_lc_ibsubmit(prs):
    s = blank(prs)
    header(s, "Lifecycle 2 & 3 — DB Claim → IB Submit",
           "Broker polls DB, claims commands, builds OCO brackets")
    lifecycle_bar(s, 3)

    box(s, 0.4, 1.4, 5.7, 5.4, fill=PANEL, line=BLUE)
    tb(s, "Broker Submission Flow", 0.6, 1.5, 5.3, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        ("Step 1 — Claim (SUBMITTING state)", BLUE, True, 14),
        "  Write status=SUBMITTING BEFORE calling IB",
        "  → crash-safe: restart won't duplicate orders",
        "",
        ("Step 2 — Build OCO bracket", BLUE, True, 14),
        "  parent  = LimitOrder('BUY', 1, 6250)",
        "  tp      = LimitOrder('SELL',1, 6252)",
        "  sl      = StopOrder( 'SELL',1, 6248)",
        "  sl.transmit = True  ← sends all 3",
        "",
        ("Step 3 — Submit to PAPER port 4002", BLUE, True, 14),
        "  ib_paper.placeOrder(contract, parent)",
        "  ib_paper.placeOrder(contract, tp)",
        "  ib_paper.placeOrder(contract, sl)",
        "",
        ("Step 4 — Write IB order IDs to DB", BLUE, True, 14),
        "  status = SUBMITTED",
        "  ib_parent_order_id = 1001",
        "  ib_tp_order_id     = 1002",
        "  ib_sl_order_id     = 1003",
    ], 0.6, 2.0, 5.3, 4.6, default_size=13)

    box(s, 6.3, 1.4, 6.6, 2.6, fill=PANEL, line=YELLOW)
    tb(s, "IB Status after submission", 6.5, 1.5, 6.2, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        ("orderId=1001  LMT BUY  6250  PreSubmitted", BLUE,   False, 12),
        ("orderId=1002  LMT SELL 6252  PreSubmitted", GREEN,  False, 12),
        ("orderId=1003  STP SELL 6248  PreSubmitted", RED,    False, 12),
        ("orderId=1004  LMT BUY  6250  PreSubmitted", BLUE,   False, 12),
        ("orderId=1005  LMT SELL 6254  PreSubmitted", GREEN,  False, 12),
        ("orderId=1006  STP SELL 6246  PreSubmitted", RED,    False, 12),
        ("orderId=1007  STP SELL 6250  PreSubmitted", YELLOW, False, 12),
        ("orderId=1008  LMT BUY  6248  PreSubmitted", GREEN,  False, 12),
        ("... 4 more for cmd id=4 ...",               GRAY,   False, 12),
    ], 6.5, 2.0, 6.2, 1.9, default_size=12)

    box(s, 6.3, 4.15, 6.6, 2.65, fill=PANEL, line=GOLD)
    tb(s, "State Machine — Broker path", 6.5, 4.25, 6.2, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        ("PENDING", YELLOW, True, 14), ("  →  SUBMITTING  →  SUBMITTED", WHITE, False, 14),
        "",
        "SUBMITTING is the claim lock.",
        "If Broker crashes here and restarts:",
        "  • Check IB for matching order_ref",
        "  • Found   → set SUBMITTED",
        "  • Not found → reset to PENDING, retry",
    ], 6.5, 4.7, 6.2, 2.0, default_size=13)

def slide_lc_fill(prs):
    s = blank(prs)
    header(s, "Lifecycle 4 — Price Hits 6250",
           "IB fires callbacks  →  Broker updates DB  →  OCO children activate")
    lifecycle_bar(s, 4)

    box(s, 0.4, 1.4, 4.0, 5.4, fill=PANEL, line=GREEN)
    tb(s, "10:23 CT — Price = 6250", 0.6, 1.5, 3.6, 0.4, size=15, bold=True, color=GREEN)
    mtb(s, [
        "IB evaluates all open orders:",
        "",
        ("LMT BUY  1001 @ 6250  →  FILLS", GREEN, True, 14),
        ("LMT BUY  1004 @ 6250  →  FILLS", GREEN, True, 14),
        ("STP SELL 1007 @ 6250  →  FILLS", GREEN, True, 14),
        ("STP SELL 1010 @ 6250  →  FILLS", GREEN, True, 14),
        "",
        "All 4 parents fill simultaneously.",
        "",
        ("Net IB position = 0", YELLOW, True, 15),
        "2 LONG + 2 SHORT = flat",
        "",
        ("The system tracks 4 independent", BLUE, False, 13),
        ("virtual strategy legs.", BLUE, False, 13),
        ("This is intentional.", BLUE, True, 13),
    ], 0.6, 2.0, 3.6, 4.6, default_size=14)

    box(s, 4.6, 1.4, 4.2, 5.4, fill=PANEL, line=GOLD)
    tb(s, "IB Callbacks to Broker", 4.8, 1.5, 3.8, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        ("orderStatus:", GOLD, True, 13),
        "  orderId=1001 Filled @ 6250",
        "  orderId=1004 Filled @ 6250",
        "  orderId=1007 Filled @ 6250",
        "  orderId=1010 Filled @ 6250",
        "",
        ("Children now active:", GOLD, True, 13),
        "  orderId=1002  Submitted (TP)",
        "  orderId=1003  Submitted (SL)",
        "  orderId=1005  Submitted (TP)",
        "  orderId=1006  Submitted (SL)",
        "  ...and 4 more",
    ], 4.8, 2.0, 3.8, 4.5, default_size=13)

    box(s, 9.0, 1.4, 4.0, 5.4, fill=PANEL, line=BLUE)
    tb(s, "DB After Fill", 9.2, 1.5, 3.6, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        ("UPDATE commands SET", BLUE, True, 12),
        "  status     = 'FILLED'",
        "  fill_price = 6250.00",
        "  filled_at  = '10:23:14'",
        "WHERE id IN (1,2,3,4);",
        "",
        ("INSERT INTO positions", BLUE, True, 12),
        "  (command_id, symbol,",
        "   direction, entry_price,",
        "   entry_time)",
        "  ×4 rows",
        "",
        ("INSERT INTO ib_events", BLUE, True, 12),
        "  fill × 4 entries",
    ], 9.2, 2.0, 3.6, 4.5, default_size=12)

def slide_lc_exit(prs):
    s = blank(prs)
    header(s, "Lifecycle 5 & 6 — OCO Active + Exit Scenarios",
           "Three ways a position closes — each triggers replenishment logic")
    lifecycle_bar(s, 6)

    # Scenario A
    box(s, 0.3, 1.4, 4.1, 5.5, fill=PANEL, line=GREEN)
    tb(s, "Scenario A — TP Hit", 0.5, 1.5, 3.7, 0.4, size=15, bold=True, color=GREEN)
    mtb(s, [
        "Price rises to 6252.",
        "",
        ("LMT SELL 1002 @ 6252  FILLS", GREEN, True, 13),
        ("STP SELL 1003 @ 6248  CANCELLED", RED,   True, 13),
        "(OCO — one cancels other)",
        "",
        "DB update:",
        "  exit_price  = 6252.00",
        "  exit_reason = 'TP'",
        "  pnl_points  = +2.0",
        "  status      = CLOSED",
        "",
        ("→ Replenishment: IMMEDIATE", GREEN, True, 14),
        "  No cool-down. New command",
        "  written to DB right away.",
    ], 0.5, 2.0, 3.7, 4.7, default_size=13)

    # Scenario B
    box(s, 4.6, 1.4, 4.1, 5.5, fill=PANEL, line=RED)
    tb(s, "Scenario B — SL Hit", 4.8, 1.5, 3.7, 0.4, size=15, bold=True, color=RED)
    mtb(s, [
        "Price drops to 6248.",
        "",
        ("STP SELL 1003 @ 6248  FILLS",   RED,   True, 13),
        ("LMT SELL 1002 @ 6252  CANCELLED",GREEN, True, 13),
        "",
        "DB update:",
        "  exit_price  = 6248.00",
        "  exit_reason = 'SL'",
        "  pnl_points  = -2.0",
        "  status      = CLOSED",
        "",
        ("→ Cool-down: 30 seconds", RED, True, 14),
        "  Line disarmed for 30s.",
        "  Avoids whipsaw re-entry.",
        "  Replenishment after cool-down.",
    ], 4.8, 2.0, 3.7, 4.7, default_size=13)

    # Scenario C
    box(s, 8.9, 1.4, 4.1, 5.5, fill=PANEL, line=YELLOW)
    tb(s, "Scenario C — Stagnation", 9.1, 1.5, 3.7, 0.4, size=15, bold=True, color=YELLOW)
    mtb(s, [
        "Broker monitors every open position.",
        "",
        "Check every few seconds:",
        "  time_in_trade > 60s",
        "  AND",
        "  |current - entry| < 0.5pt",
        "",
        ("→ Dead line detected", YELLOW, True, 14),
        "",
        "Action:",
        "  Market exit order",
        "  Cancel TP and SL",
        "",
        "  exit_reason = 'STAGNATION'",
        "  pnl ≈ 0 (slippage only)",
        "",
        ("→ Replenishment: IMMEDIATE", GREEN, True, 13),
    ], 9.1, 2.0, 3.7, 4.7, default_size=13)

def slide_lc_replenish(prs):
    s = blank(prs)
    header(s, "Lifecycle 7 — Replenishment",
           "Decider detects closed commands and re-arms the line")
    lifecycle_bar(s, 7)

    box(s, 0.4, 1.4, 5.8, 5.4, fill=PANEL, line=GOLD)
    tb(s, "Decider Replenishment Loop", 0.6, 1.5, 5.4, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        "Runs every 10s. Queries DB:",
        "",
        ("SELECT * FROM commands", BLUE, True, 13),
        "WHERE status = 'CLOSED'",
        "AND replenishment_issued = 0",
        "AND session_state != 'SHUTDOWN'",
        "",
        "For each result:",
        "  1. Check cooldown (SL exits only)",
        "  2. Fetch current price",
        ("  3. Re-evaluate toggle", GOLD, True, 14),
        "     current > line  → LMT BUY + STP SELL",
        "     current < line  → STP BUY + LMT SELL",
        "  4. Write new PENDING command",
        "  5. Set replenishment_issued = 1",
        "     on original command",
        "",
        ("Invariant: exactly ONE replenishment", YELLOW, True, 13),
        ("per closed command. Never zero, never two.", YELLOW, False, 13),
    ], 0.6, 2.0, 5.4, 4.6, default_size=13)

    box(s, 6.5, 1.4, 6.4, 5.4, fill=PANEL, line=BLUE)
    tb(s, "Cool-down Tracking (SL only)", 6.7, 1.5, 6.0, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        "Cool-down is side-specific:",
        "",
        ("system_state table:", BLUE, True, 13),
        "  key   = 'cooldown_MES_6250_SELL'",
        "  value = '10:24:52'  (expires at)",
        "",
        "BUY and SELL sides tracked independently.",
        "SL on BUY side does not delay SELL side.",
        "",
        ("Shutdown overrides everything:", RED, True, 14),
        "Once SHUTDOWN is in system_state,",
        "replenishment query returns 0 rows.",
        "No new commands. Ever.",
        "",
        ("Replenishment = the engine.", GOLD, True, 15),
        "It's why the system keeps harvesting",
        "20 hits per day instead of just 1.",
    ], 6.7, 2.0, 6.0, 4.6, default_size=13)

def slide_lc_shutdown(prs):
    s = blank(prs)
    header(s, "Lifecycle 8 — Shutdown Sequence",
           "T-60 minutes before CME close — orderly then panic")
    lifecycle_bar(s, 8)

    steps = [
        ("Step 1\nT-60min", "SHUTDOWN written\nto system_state.\nReplenishment = DEAD.", RED),
        ("Step 2\nT-60min", "Cancel ALL\nSUBMITTED orders\n(unfilled entries).", YELLOW),
        ("Step 3\nT-60min", "Tighten all open\nbracket stops\nto 1 point.", YELLOW),
        ("Step 4\nOrderly", "Exit open positions\none by one.\n30s patience each.", GOLD),
        ("Step 5\nPanic", "T-10min: market\nexit ALL remaining\nsimultaneously.", RED),
    ]

    bx = 0.4
    bw = 2.4
    by = 1.5
    bh = 3.0
    for i, (title, desc, color) in enumerate(steps):
        x = bx + i * (bw + 0.16)
        box(s, x, by, bw, bh, fill=PANEL, line=color)
        tb(s, title, x+0.1, by+0.1, bw-0.2, 0.65, size=14, bold=True, color=color,
           align=PP_ALIGN.CENTER)
        hline(s, by+0.8, color, x=x+0.1, w=bw-0.2, thickness=1)
        tb(s, desc, x+0.1, by+0.9, bw-0.2, bh-1.1,
           size=13, color=GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            tb(s, "→", x+bw+0.0, by+1.2, 0.2, 0.4,
               size=20, color=RED, bold=True, align=PP_ALIGN.CENTER)

    mtb(s, [
        ("Exit reasons logged to DB:", GOLD, True, 14),
        ("  SHUTDOWN_ORDERLY  ", YELLOW, False, 14), (" (step 4)  |  ", GRAY, False, 14),
        ("SHUTDOWN_PANIC",RED, False, 14), (" (step 5)", GRAY, False, 14),
    ], 0.5, 4.7, 12.33, 0.5, default_size=14)

    tbl(s,
        ["Config key", "Default", "Purpose"],
        [
            ["session.shutdown_offset_minutes",  "60",  "When shutdown begins (before close)"],
            ["shutdown.exit_patience_seconds",   "30",  "Wait between orderly exits"],
            ["shutdown.panic_threshold_minutes", "10",  "Switch to panic mode"],
        ],
        0.4, 5.3, 12.33, 1.6
    )

def slide_state_machine(prs):
    s = blank(prs)
    header(s, "Command Status State Machine",
           "All allowed transitions. No skipping states.")

    box(s, 0.4, 1.4, 5.0, 5.5, fill=PANEL, line=GOLD)
    tb(s, "States", 0.6, 1.5, 4.6, 0.4, size=15, bold=True, color=GOLD)
    tbl(s,
        ["Status", "Set by", "Meaning"],
        [
            [("PENDING",    YELLOW), "Decider", "Written, waiting for Broker"],
            [("SUBMITTING", BLUE),   "Broker",  "Claim lock — before IB call"],
            [("SUBMITTED",  BLUE),   "Broker",  "In IB, waiting for fill"],
            [("FILLED",     GREEN),  "Broker",  "Entry filled, children active"],
            [("EXITING",    YELLOW), "Broker",  "Market exit in progress"],
            [("CLOSED",     GREEN),  "Broker",  "Fully resolved. Triggers replenish."],
            [("CANCELLED",  GRAY),   "Broker",  "Cancelled before fill"],
            [("ERROR",      RED),    "Broker",  "IB rejected. Manual review."],
            [("RECONCILE\nREQUIRED",RED),"Broker","DB and IB disagree. Investigate."],
        ],
        0.5, 2.0, 4.8, 4.7
    )

    box(s, 5.7, 1.4, 7.3, 5.5, fill=PANEL, line=BLUE)
    tb(s, "Transition Rules", 5.9, 1.5, 7.0, 0.4, size=15, bold=True, color=GOLD)
    mtb(s, [
        ("PENDING",     YELLOW, True,  15), " → SUBMITTING",
        ("SUBMITTING",  BLUE,   True,  15), " → SUBMITTED | ERROR",
        ("SUBMITTED",   BLUE,   True,  15), " → FILLED | CANCELLED",
        ("FILLED",      GREEN,  True,  15), " → EXITING",
        ("EXITING",     YELLOW, True,  15), " → CLOSED",
        ("CLOSED",      GREEN,  True,  15), " (terminal — new replenishment command written)",
        ("CANCELLED",   GRAY,   True,  15), " (terminal)",
        ("ERROR",       RED,    True,  15), " (terminal — manual review required)",
        ("RECONCILE",   RED,    True,  15), " → SUBMITTED | FILLED | CANCELLED | ERROR",
        "",
        ("No other transitions valid.", RED, True, 14),
        "Skipped states must be logged as anomalies.",
        "",
        ("Failure rules:", GOLD, True, 14),
        "•  SUBMITTING on restart → check IB → set SUBMITTED or PENDING",
        "•  Duplicate fill callback → ignore, log DEBUG",
        "•  DB write fails during update → log to file, abort",
        "•  IB disconnect → 5 retries over 5 min → abort session",
    ], 5.9, 2.0, 7.0, 4.6, default_size=13)

def slide_architecture(prs):
    s = blank(prs)
    header(s, "Architecture", "5 components communicating through a single SQLite database")

    # Component boxes
    comps = [
        (0.4,  2.0, "DECIDER",    "Generates commands\nHandles replenishment\nTriggers shutdown",      GOLD),
        (3.5,  2.0, "BROKER",     "Executes IB orders\nUpdates DB status\nMonitors stagnation",        BLUE),
        (6.6,  2.0, "FETCHER",    "Downloads history CSV\nLIVE port only\nRuns independently",         GRAY),
        (9.7,  2.0, "VISUALIZER", "Browser dashboard\nDB viewer\nLive P&L + status",                  GREEN),
        (3.5,  5.0, "ANALYZER",   "A/B test results\nP&L by bracket + strength\n[LATER]",             GRAY),
    ]
    for x, y, title, desc, color in comps:
        box(s, x, y, 2.8, 2.4, fill=PANEL, line=color)
        tb(s, title, x+0.1, y+0.1, 2.6, 0.5, size=15, bold=True, color=color,
           align=PP_ALIGN.CENTER)
        hline(s, y+0.65, color, x=x+0.1, w=2.6, thickness=1)
        tb(s, desc, x+0.1, y+0.75, 2.6, 1.5, size=12, color=GRAY,
           align=PP_ALIGN.CENTER)

    # DB box in center
    box(s, 5.2, 3.5, 2.8, 1.5, fill=RGBColor(0x30,0x45,0x60), line=GOLD)
    tb(s, "SQLite DB\ngalao.db", 5.3, 3.6, 2.6, 1.2, size=16, bold=True,
       color=GOLD, align=PP_ALIGN.CENTER)

    # IB Gateway box
    box(s, 9.7, 4.8, 2.8, 1.8, fill=PANEL, line=BLUE)
    tb(s, "IB GATEWAY", 9.8, 4.9, 2.6, 0.4, size=13, bold=True, color=BLUE,
       align=PP_ALIGN.CENTER)
    mtb(s, [
        ("LIVE  port 4001", GRAY,  False, 12),
        ("PAPER port 4002", GREEN, False, 12),
    ], 9.8, 5.4, 2.6, 0.9, default_size=12, align=PP_ALIGN.CENTER)

    # Rule box
    box(s, 0.4, 5.5, 2.8, 1.5, fill=PANEL, line=YELLOW)
    mtb(s, [
        ("One rule:", YELLOW, True, 14),
        "Components talk",
        "through DB only.",
        "No direct calls.",
    ], 0.5, 5.6, 2.6, 1.2, default_size=13, align=PP_ALIGN.CENTER)

def slide_schema(prs):
    s = blank(prs)
    header(s, "Database Schema — commands table",
           "Primary table. Every order's full lifecycle lives here.")

    tbl(s,
        ["Column", "Type", "Description"],
        [
            ["id",                   "INT PK",  "Auto increment"],
            [("status",GOLD),        "TEXT",    "PENDING→SUBMITTING→SUBMITTED→FILLED→EXITING→CLOSED"],
            [("direction",BLUE),     "TEXT",    "BUY / SELL"],
            [("entry_order_type",BLUE),"TEXT",  "LMT or STP  (toggle result)"],
            ["symbol",               "TEXT",    "e.g. MES"],
            ["contract_month",       "TEXT",    "e.g. MESM6"],
            ["entry_price",          "REAL",    "Critical line price (rounded to 0.25)"],
            [("bracket_size",GOLD),  "REAL",    "A/B test variable — e.g. 2.0 or 4.0"],
            ["take_profit",          "REAL",    "entry ± bracket_size"],
            ["stop_loss",            "REAL",    "entry ∓ bracket_size"],
            [("line_strength",GOLD), "INT",     "A/B test variable — 1 / 2 / 3"],
            ["ib_parent_order_id",   "INT",     "IB order ID for entry"],
            ["ib_tp_order_id",       "INT",     "IB order ID for TP"],
            ["ib_sl_order_id",       "INT",     "IB order ID for SL"],
            ["fill_price / filled_at","REAL/DT","Actual fill details"],
            [("exit_reason",YELLOW), "TEXT",    "TP / SL / STAGNATION / SHUTDOWN_ORDERLY / SHUTDOWN_PANIC"],
            [("replenishment_issued",GREEN),"BOOL","Prevents double replenishment (invariant I-03)"],
            ["claimed_at",           "DT",      "When SUBMITTING was set — crash recovery anchor"],
            ["is_replenishment",     "BOOL",    "True if this is a replenishment command"],
            ["parent_command_id",    "INT",     "FK to original command that was filled"],
            ["order_ref",            "TEXT",    "Stored in IB orderRef field — used for reconciliation"],
        ],
        0.3, 1.4, 12.7, 5.5
    )

def slide_schema2(prs):
    s = blank(prs)
    header(s, "Database Schema — Supporting Tables")

    box(s, 0.3, 1.4, 4.0, 5.5, fill=PANEL, line=BLUE)
    tb(s, "positions", 0.5, 1.5, 3.6, 0.4, size=15, bold=True, color=GOLD)
    tbl(s,
        ["Column", "Type"],
        [
            ["id",          "INT PK"],
            ["command_id",  "INT FK"],
            ["symbol",      "TEXT"],
            ["direction",   "TEXT"],
            ["entry_price", "REAL"],
            ["entry_time",  "DT"],
            ["exit_price",  "REAL"],
            ["exit_time",   "DT"],
            [("pnl_points", GREEN), "REAL"],
            [("exit_reason",YELLOW),"TEXT"],
        ],
        0.4, 2.0, 3.8, 4.7
    )
    tb(s, "← P&L source of truth", 0.5, 6.75, 3.5, 0.3, size=11, color=GOLD, italic=True)

    box(s, 4.5, 1.4, 4.2, 2.6, fill=PANEL, line=YELLOW)
    tb(s, "ib_events", 4.7, 1.5, 3.8, 0.4, size=15, bold=True, color=GOLD)
    tbl(s,
        ["Column", "Type"],
        [
            ["id",          "INT PK"],
            ["timestamp",   "DT UTC"],
            ["event_type",  "TEXT"],
            ["ib_order_id", "INT"],
            ["data",        "JSON"],
        ],
        4.6, 2.0, 4.0, 1.8
    )

    box(s, 4.5, 4.15, 4.2, 2.75, fill=PANEL, line=GREEN)
    tb(s, "system_state", 4.7, 4.25, 3.8, 0.4, size=15, bold=True, color=GOLD)
    tbl(s,
        ["key (TEXT PK)", "value"],
        [
            ["session_state",        ("RUNNING / SHUTDOWN", RED)],
            ["cooldown_MES_6250_BUY","expiry timestamp"],
            ["cooldown_MES_6250_SELL","expiry timestamp"],
        ],
        4.6, 4.7, 4.0, 1.8
    )

    box(s, 8.9, 1.4, 4.2, 5.5, fill=PANEL, line=GRAY)
    tb(s, "critical_lines  &  release_notes", 9.1, 1.5, 3.8, 0.4, size=14, bold=True, color=GOLD)
    tbl(s,
        ["Table", "Key Columns"],
        [
            ["critical_lines", "date, symbol, price, strength, source"],
            ["release_notes",  "timestamp, program, version, change_type, description"],
        ],
        9.0, 2.0, 4.0, 1.3
    )
    mtb(s, [
        ("release_notes.py --program <name>", BLUE, True, 13),
        "Filter release notes by component.",
        "",
        ("test_galao.db", YELLOW, True, 13),
        "Separate DB for regression tests.",
        "Never touches production DB.",
    ], 9.1, 3.5, 3.8, 3.2, default_size=13)

def slide_abtesting(prs):
    s = blank(prs)
    header(s, "A/B Testing Framework",
           "The system is a learning platform — every trade is a data point")

    box(s, 0.4, 1.4, 5.8, 5.5, fill=PANEL, line=GOLD)
    tb(s, "Two Test Variables", 0.6, 1.5, 5.4, 0.4, size=16, bold=True, color=GOLD)
    tbl(s,
        ["Variable", "Values", "Where set"],
        [
            [("bracket_size", GOLD), "[2, 4, 8, 16] pts\n(any positive value)", "config.yaml"],
            [("line_strength",BLUE), "1 / 2 / 3",           "levels_daily file"],
        ],
        0.5, 2.0, 5.6, 1.8
    )
    mtb(s, [
        "",
        "Both recorded on every command row.",
        "No runtime decisions based on them.",
        "Collection only in V1.",
        "",
        ("Future Analyzer will answer:", GOLD, True, 15),
        "  •  Which bracket size has highest win rate?",
        "  •  Do strength=3 lines outperform strength=1?",
        "  •  What is win rate × bracket size?",
        "  •  P&L per symbol, per day",
    ], 0.6, 3.9, 5.4, 2.8, default_size=14)

    box(s, 6.5, 1.4, 6.4, 5.5, fill=PANEL, line=BLUE)
    tb(s, "Gemini Suggested Batches", 6.7, 1.5, 6.0, 0.4, size=16, bold=True, color=GOLD)
    tbl(s,
        ["Batch", "Size", "Profile"],
        [
            [("A — Micro",   GREEN),  "1.5 / 1.5 pt", "Highest freq, lowest profit/trade"],
            [("B — Standard",BLUE),   "3.0 / 3.0 pt", "Balanced"],
            [("C — Wide",    YELLOW), "6.0 / 6.0 pt", "Lower freq, higher profit"],
            [("D — Runner",  RED),    "12.0 / 12 pt", "Rare, high value"],
        ],
        6.6, 2.0, 6.2, 2.3
    )
    mtb(s, [
        "",
        ("Success metric:", GOLD, True, 14),
        "  Total P&L ÷ Number of fills",
        "  → identifies the 'sweet spot'",
        "",
        "Run 5 trades per batch to get",
        "initial signal on each bracket size.",
        "",
        ("Hypothesis:", BLUE, True, 14),
        "Small brackets (2pt) win more often",
        "because the line's 4× traffic means",
        "price returns before reaching TP on",
        "larger brackets.",
    ], 6.7, 4.4, 6.0, 2.3, default_size=13)

def slide_devquality(prs):
    s = blank(prs)
    header(s, "Dev Quality Framework",
           "Self-test, regression, versioning — built in from day one")

    items = [
        ("--self-test flag", GOLD,
         ["Every Python script has --self-test",
          "Exit 0 = pass,  Exit 1 = fail",
          "GUI programs test headless (no browser)",
          "Task = complete ONLY after self-test passes"]),
        ("Versioning Workflow", BLUE,
         ["1. Copy file → versions/{name}.{YYYYMMDD_HHMM}",
          "2. Make changes",
          "3. Run --self-test → must pass",
          "4. Write to release_notes DB table",
          "5. Announce complete"]),
        ("Regression Tests", GREEN,
         ["regression.py  (on demand only)",
          "Layer 1: --self-test all components",
          "Layer 2: feature/logic tests (no IB)",
          "Layer 3: real IB — submit + cancel order",
          "--quick  |  --layer3-only  |  --program <x>"]),
        ("Pre-flight Checklist", YELLOW,
         ["Runs before every trading session",
          "1. LIVE port 4001 connection",
          "2. PAPER port 4002 connection",
          "3. Price fetch from LIVE",
          "4. DB read/write test",
          "Any failure = hard abort"]),
    ]

    x = 0.35
    for i, (title, color, pts) in enumerate(items):
        bx = x + i * 3.25
        box(s, bx, 1.4, 3.0, 5.5, fill=PANEL, line=color)
        tb(s, title, bx+0.1, 1.5, 2.8, 0.5, size=14, bold=True, color=color,
           align=PP_ALIGN.CENTER)
        hline(s, 2.1, color, x=bx+0.1, w=2.8, thickness=1)
        lines = [(p, GRAY, False, 13) for p in pts]
        mtb(s, lines, bx+0.1, 2.2, 2.8, 4.5, default_size=13)

def slide_next_steps(prs):
    s = blank(prs)
    header(s, "Next Steps")

    box(s, 0.4, 1.4, 5.8, 5.5, fill=PANEL, line=GOLD)
    tb(s, "Build Order (V1)", 0.6, 1.5, 5.4, 0.4, size=16, bold=True, color=GOLD)
    mtb(s, [
        ("Phase 1 — Foundation", GOLD, True, 15),
        "  lib/config_loader.py",
        "  lib/db.py  (schema init)",
        "  lib/logger.py",
        "  preflight.py",
        "",
        ("Phase 2 — Fetcher", BLUE, True, 15),
        "  fetcher.py  (from V1 code)",
        "  Verify CSV output format",
        "",
        ("Phase 3 — Broker", BLUE, True, 15),
        "  lib/ib_client.py",
        "  lib/order_builder.py",
        "  broker.py",
        "",
        ("Phase 4 — Decider", GOLD, True, 15),
        "  decider.py  (toggle + replenishment)",
        "",
        ("Phase 5 — Visualizer", GREEN, True, 15),
        "  DB viewer + live dashboard",
    ], 0.6, 2.0, 5.4, 4.6, default_size=13)

    box(s, 6.5, 1.4, 6.4, 2.5, fill=PANEL, line=GREEN)
    tb(s, "Design Documents", 6.7, 1.5, 6.0, 0.4, size=16, bold=True, color=GOLD)
    mtb(s, [
        ("rules_book.md       ", BLUE, True, 14), "14 sections, all rules",
        ("design_book.md      ", BLUE, True, 14), "full architecture + schema",
        ("tech_solutions_book.md", BLUE, True, 14), "IB ports, code patterns",
        ("walkthrough_book.md ", BLUE, True, 14), "complete trade lifecycle trace",
        ("release_notes.md    ", BLUE, True, 14), "all decisions logged",
    ], 6.7, 2.0, 6.0, 1.7, default_size=13)

    box(s, 6.5, 4.05, 6.4, 2.85, fill=PANEL, line=YELLOW)
    tb(s, "Key Design Decisions", 6.7, 4.15, 6.0, 0.4, size=16, bold=True, color=GOLD)
    mtb(s, [
        ("✓", GREEN, True, 14),  " PAPER port only until approved",
        ("✓", GREEN, True, 14),  " Virtual strategy legs (not broker net)",
        ("✓", GREEN, True, 14),  " LMT/STP toggle per side",
        ("✓", GREEN, True, 14),  " Stagnation kill-switch (60s / 0.5pt)",
        ("✓", GREEN, True, 14),  " Replenishment fully dead on shutdown",
        ("✓", GREEN, True, 14),  " Single broker process",
        ("✓", GREEN, True, 14),  " Partial fills ignored in V1",
    ], 6.7, 4.6, 6.0, 2.1, default_size=13)

# ── Main ──────────────────────────────────────────────────────────────────────

def build_presentation():
    prs = new_prs()
    slide_title_page(prs)
    slide_agenda(prs)
    slide_thesis(prs)
    slide_strategy(prs)
    slide_critical_lines(prs)
    slide_toggle(prs)
    slide_bracket(prs)
    slide_lifecycle_overview(prs)
    slide_lc_daystart(prs)
    slide_lc_ibsubmit(prs)
    slide_lc_fill(prs)
    slide_lc_exit(prs)
    slide_lc_replenish(prs)
    slide_lc_shutdown(prs)
    slide_state_machine(prs)
    slide_architecture(prs)
    slide_schema(prs)
    slide_schema2(prs)
    slide_abtesting(prs)
    slide_devquality(prs)
    slide_next_steps(prs)
    out = "galao_presentation.pptx"
    prs.save(out)
    print(f"[OK] Saved: {out}  ({prs.slides.__len__()} slides)")
    return out

def self_test():
    try:
        from pptx import Presentation as P
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        prs = new_prs()
        s = blank(prs)
        tb(s, "test", 0, 0, 1, 0.5)
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            tmp = f.name
        prs.save(tmp)
        os.unlink(tmp)
        print("[self-test] PASS")
        return True
    except Exception as e:
        print(f"[self-test] FAIL: {e}")
        return False

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--self-test", action="store_true")
    args = parser.parse_args()
    if args.self_test:
        sys.exit(0 if self_test() else 1)
    build_presentation()
